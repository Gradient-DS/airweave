"""Simple PPTX text extraction using python-pptx."""

import asyncio
from pathlib import Path
from typing import Dict, List

from pptx import Presentation

from airweave.platform.sync.async_helpers import run_in_thread_pool
from airweave.platform.sync.exceptions import EntityProcessingError
from airweave.platform.converters._base import BaseTextConverter


class SimplePptxConverter(BaseTextConverter):
    """Extract text from PPTX files using python-pptx (local, no API)."""

    BATCH_SIZE = 10
    MAX_CONCURRENT = 20

    def __init__(self):
        """Initialize SimplePptxConverter."""
        super().__init__()
        self.logger.info("Initialized SimplePptxConverter (local extraction)")

    async def convert_batch(self, paths: List[Path]) -> Dict[Path, str]:
        """Extract text from multiple PPTX files.

        Args:
            paths: List of PPTX file paths

        Returns:
            Dict mapping path to extracted markdown text
        """
        self.logger.debug(f"Converting {len(paths)} PPTX files")

        semaphore = asyncio.Semaphore(self.MAX_CONCURRENT)
        tasks = [self._convert_single(path, semaphore) for path in paths]
        results = await asyncio.gather(*tasks, return_exceptions=True)

        output = {}
        for path, result in zip(paths, results):
            if isinstance(result, Exception):
                self.logger.error(f"Failed to convert {path.name}: {result}")
                output[path] = None
            else:
                output[path] = result

        success_count = sum(1 for v in output.values() if v is not None)
        self.logger.info(
            f"Converted {success_count}/{len(paths)} PPTX files successfully"
        )

        return output

    async def _convert_single(
        self, path: Path, semaphore: asyncio.Semaphore
    ) -> str:
        """Extract text from a single PPTX file.

        Args:
            path: Path to PPTX file
            semaphore: Concurrency control

        Returns:
            Extracted text as markdown
        """
        async with semaphore:
            try:
                def extract():
                    prs = Presentation(str(path))
                    slides_text = []

                    for slide_num, slide in enumerate(prs.slides, 1):
                        slide_content = [f"## Slide {slide_num}"]

                        # Extract text from shapes
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text:
                                text = shape.text.strip()
                                if text:
                                    slide_content.append(text)

                            # Extract table content
                            if shape.has_table:
                                table = shape.table
                                for row in table.rows:
                                    row_text = " | ".join(
                                        cell.text.strip() for cell in row.cells
                                    )
                                    slide_content.append(row_text)

                        if len(slide_content) > 1:  # Has content beyond slide number
                            slides_text.append("\n\n".join(slide_content))

                    return "\n\n---\n\n".join(slides_text)

                text = await run_in_thread_pool(extract)

                if not text or not text.strip():
                    self.logger.warning(f"No text extracted from {path.name}")
                    return ""

                return text.strip()

            except Exception as e:
                self.logger.error(
                    f"Error extracting text from {path.name}: {e}",
                    exc_info=True
                )
                raise EntityProcessingError(
                    f"Failed to extract text from PPTX: {e}"
                )
